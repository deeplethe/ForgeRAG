"""
Document format conversion: Office/HTML/Text → PDF.

Pure Python implementation — no external software dependencies.

Supported inputs:
    - DOCX (.docx)     → python-docx + fpdf2
    - PPTX (.pptx)     → python-pptx + fpdf2
    - HTML (.html/.htm) → beautifulsoup4 + fpdf2
    - Markdown (.md)    → markdown + fpdf2
    - Plain text (.txt) → fpdf2
    - XLSX (.xlsx)      → openpyxl + fpdf2

The converted PDF preserves text structure (headings, paragraphs,
tables, slides) and is used for:
    - Parsing (full bbox + page coordinates via PyMuPDF)
    - Frontend PDF viewer (preview + highlight)

Required pip packages:
    pip install fpdf2
    pip install python-docx    # for DOCX
    pip install python-pptx    # for PPTX
    pip install beautifulsoup4 # for HTML
    pip install markdown        # for MD
    pip install openpyxl        # for XLSX
"""

from __future__ import annotations

import contextlib
import logging
import re
from pathlib import Path

log = logging.getLogger(__name__)

# Extensions that need conversion to PDF before parsing.
#
# Office formats: only the OOXML variants (``.docx`` / ``.pptx`` /
# ``.xlsx``) are supported. The legacy binary formats (``.doc`` /
# ``.ppt`` / ``.xls``) are NOT here because ``python-docx`` /
# ``python-pptx`` / ``openpyxl`` are all OOXML-only — they cannot
# read the OLE Compound Document binary format. Earlier versions
# pretended to handle them with the comment "best-effort; may not
# parse" but the converter would always raise ``zipfile.BadZipFile``
# at import time. The upload route now rejects them with HTTP 415
# and a clear "save as .docx / .pptx / .xlsx" message instead of
# letting the ingest fail confusingly mid-pipeline.
CONVERTIBLE_EXTENSIONS = {
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".md",
    ".txt",
}

# Legacy Office binary formats — explicitly unsupported. Surfaced as
# its own constant so the upload route can match against it for a
# targeted error message rather than the generic "unsupported file".
LEGACY_OFFICE_EXTENSIONS = (".doc", ".ppt", ".xls")


def needs_conversion(path: str | Path) -> bool:
    """Return True if this file should be converted to PDF before parsing."""
    return Path(path).suffix.lower() in CONVERTIBLE_EXTENSIONS


def convert_to_pdf(
    source_path: str | Path,
    *,
    output_dir: str | Path | None = None,
) -> Path:
    """
    Convert a document to PDF using pure Python libraries.

    Returns the path to the generated PDF file.
    Raises RuntimeError if conversion fails or required library is missing.
    """
    source = Path(source_path)
    if not source.exists():
        raise FileNotFoundError(f"source file not found: {source}")

    if output_dir is None:
        output_dir = source.parent
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    ext = source.suffix.lower()
    out_path = output_dir / (source.stem + ".pdf")

    converters = {
        ".docx": _convert_docx,
        ".pptx": _convert_pptx,
        ".xlsx": _convert_xlsx,
        ".html": _convert_html,
        ".htm": _convert_html,
        ".md": _convert_markdown,
        ".txt": _convert_text,
    }

    converter = converters.get(ext)
    if converter is None:
        raise RuntimeError(f"no converter for extension: {ext}")

    log.info("converting %s → PDF", source.name)
    converter(source, out_path)
    log.info("conversion complete: %s (%d bytes)", out_path.name, out_path.stat().st_size)
    return out_path


# ============================================================================
# PDF builder helpers
# ============================================================================


def _make_pdf():
    """Create an FPDF instance with CJK font support."""
    try:
        from fpdf import FPDF
    except ImportError:
        raise RuntimeError("fpdf2 is required for document conversion: pip install fpdf2")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Try to load a CJK-capable font for Chinese/Japanese/Korean text
    font_name = _setup_cjk_font(pdf)
    if font_name:
        pdf.set_font(font_name, size=10)
    else:
        pdf.set_font("Helvetica", size=10)

    return pdf, font_name


def _setup_cjk_font(pdf) -> str | None:
    """Try to register a system CJK font. Returns font name or None."""
    candidates = [
        # Windows
        (r"C:\Windows\Fonts\msyh.ttc", "msyh"),
        (r"C:\Windows\Fonts\simsun.ttc", "simsun"),
        # Linux
        ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "noto"),
        ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "noto"),
        ("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc", "noto"),
        # macOS
        ("/System/Library/Fonts/PingFang.ttc", "pingfang"),
    ]
    for font_path, name in candidates:
        if Path(font_path).exists():
            try:
                pdf.add_font(name, "", font_path, uni=True)
                return name
            except Exception:
                continue
    return None


def _set_font(pdf, font_name: str | None, style: str = "", size: int = 10):
    """Set font, falling back to Helvetica if CJK font not available."""
    if font_name:
        # fpdf2 CJK fonts don't support B/I style on add_font, use size only
        pdf.set_font(font_name, size=size)
    else:
        pdf.set_font("Helvetica", style=style, size=size)


# ============================================================================
# DOCX converter
# ============================================================================


def _convert_docx(source: Path, out_path: Path):
    """Convert DOCX to PDF using python-docx + fpdf2."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx is required for DOCX conversion: pip install python-docx")

    doc = Document(str(source))
    pdf, font_name = _make_pdf()
    pdf.add_page()

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            _docx_render_paragraph(pdf, font_name, element, doc)
        elif tag == "tbl":
            _docx_render_table(pdf, font_name, element, doc)

    pdf.output(str(out_path))


def _docx_render_paragraph(pdf, font_name, element, doc):
    """Render a DOCX paragraph element."""
    from docx.oxml.ns import qn

    # Get paragraph style
    style_el = element.find(qn("w:pPr"))
    style_name = ""
    if style_el is not None:
        style_ref = style_el.find(qn("w:pStyle"))
        if style_ref is not None:
            style_name = style_ref.get(qn("w:val"), "")

    # Detect heading level
    outline_lvl = None
    if style_el is not None:
        ol = style_el.find(qn("w:outlineLvl"))
        if ol is not None:
            with contextlib.suppress(ValueError):
                outline_lvl = int(ol.get(qn("w:val"), "0"))

    is_heading = outline_lvl is not None or style_name.startswith("Heading")
    heading_level = 0
    if outline_lvl is not None:
        heading_level = outline_lvl
    elif style_name.startswith("Heading"):
        try:
            heading_level = int(style_name.replace("Heading", "").strip()) - 1
        except ValueError:
            heading_level = 0

    # Collect text from runs
    text = element.text or ""
    from docx.oxml.ns import qn as _qn

    for run_el in element.findall(_qn("w:r")):
        t_el = run_el.find(_qn("w:t"))
        if t_el is not None and t_el.text:
            text += t_el.text

    text = text.strip()
    if not text:
        pdf.ln(3)
        return

    # Set font based on style
    if is_heading:
        size = max(18 - heading_level * 2, 11)
        _set_font(pdf, font_name, "B", size)
        pdf.ln(4)
        _safe_multi_cell(pdf, 0, size * 0.5, text)
        pdf.ln(3)
    else:
        _set_font(pdf, font_name, "", 10)
        _safe_multi_cell(pdf, 0, 5, text)
        pdf.ln(1)


def _safe_multi_cell(pdf, w, h, text):
    """multi_cell wrapper that resets x to left margin first.

    fpdf2's multi_cell(w=0) calculates available width as
    ``page_w - right_margin - x``.  After a previous multi_cell
    the cursor x sits at the right margin, making the next call
    compute zero available width and crash.  Resetting x fixes this.
    """
    pdf.x = pdf.l_margin
    try:
        pdf.multi_cell(w, h, text)
    except Exception:
        pdf.x = pdf.l_margin
        try:
            pdf.multi_cell(w, h, "(content too wide to render)")
        except Exception:
            pass  # page is hopelessly full, skip


def _docx_render_table(pdf, font_name, element, doc):
    """Render a DOCX table element as a simple text table."""
    from docx.oxml.ns import qn

    _set_font(pdf, font_name, "", 9)
    pdf.ln(3)

    rows = element.findall(qn("w:tr"))
    for row_el in rows:
        cells = row_el.findall(qn("w:tc"))
        cell_texts = []
        for cell_el in cells:
            cell_text = ""
            for p in cell_el.findall(qn("w:p")):
                for run in p.findall(qn("w:r")):
                    t = run.find(qn("w:t"))
                    if t is not None and t.text:
                        cell_text += t.text
            cell_texts.append(cell_text.strip())
        line = " | ".join(cell_texts)
        _safe_multi_cell(pdf, 0, 4.5, line)

    pdf.ln(3)
    _set_font(pdf, font_name, "", 10)


# ============================================================================
# PPTX converter
# ============================================================================


def _convert_pptx(source: Path, out_path: Path):
    """Convert PPTX to PDF — one slide per page."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is required for PPTX conversion: pip install python-pptx")

    prs = Presentation(str(source))
    pdf, font_name = _make_pdf()

    for i, slide in enumerate(prs.slides):
        pdf.add_page()
        # Slide number header
        _set_font(pdf, font_name, "B", 8)
        pdf.cell(0, 5, f"Slide {i + 1}", ln=True)
        pdf.ln(3)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Detect title shapes
                    if shape.shape_id <= 2 and para == shape.text_frame.paragraphs[0]:
                        _set_font(pdf, font_name, "B", 14)
                        pdf.multi_cell(0, 7, text)
                        pdf.ln(3)
                    else:
                        _set_font(pdf, font_name, "", 10)
                        pdf.multi_cell(0, 5, text)
                        pdf.ln(1)

            if shape.has_table:
                _set_font(pdf, font_name, "", 9)
                pdf.ln(2)
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = " | ".join(cells)
                    _safe_multi_cell(pdf, 0, 4.5, line)
                pdf.ln(2)
                _set_font(pdf, font_name, "", 10)

    if not prs.slides:
        pdf.add_page()
        pdf.cell(0, 10, "(empty presentation)")

    pdf.output(str(out_path))


# ============================================================================
# XLSX converter
# ============================================================================


def _convert_xlsx(source: Path, out_path: Path):
    """Convert XLSX to PDF — one sheet per section."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl is required for XLSX conversion: pip install openpyxl")

    wb = load_workbook(str(source), read_only=True, data_only=True)
    pdf, font_name = _make_pdf()

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        pdf.add_page()
        _set_font(pdf, font_name, "B", 12)
        pdf.cell(0, 7, f"Sheet: {sheet_name}", ln=True)
        pdf.ln(3)
        _set_font(pdf, font_name, "", 9)

        row_count = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            _safe_multi_cell(pdf, 0, 4.5, line)
            row_count += 1
            if row_count > 2000:
                pdf.multi_cell(0, 5, f"... ({row_count}+ rows, truncated)")
                break

        if row_count == 0:
            pdf.cell(0, 5, "(empty sheet)")

    wb.close()
    pdf.output(str(out_path))


# ============================================================================
# HTML converter
# ============================================================================


def _convert_html(source: Path, out_path: Path):
    """Convert HTML to PDF by extracting text structure."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("beautifulsoup4 is required for HTML conversion: pip install beautifulsoup4")

    html = source.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")

    pdf, font_name = _make_pdf()
    pdf.add_page()

    # Walk through elements in document order
    for el in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "pre", "code"]):
        text = el.get_text(strip=True)
        if not text:
            continue

        tag = el.name
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(tag[1])
            size = max(18 - (level - 1) * 2, 11)
            _set_font(pdf, font_name, "B", size)
            pdf.ln(3)
            pdf.multi_cell(0, size * 0.5, text)
            pdf.ln(2)
        elif tag == "li":
            _set_font(pdf, font_name, "", 10)
            pdf.multi_cell(0, 5, f"• {text}")
            pdf.ln(1)
        elif tag in ("pre", "code"):
            _set_font(pdf, font_name, "", 9)
            pdf.multi_cell(0, 4.5, text)
            pdf.ln(2)
        else:
            _set_font(pdf, font_name, "", 10)
            pdf.multi_cell(0, 5, text)
            pdf.ln(1)

    pdf.output(str(out_path))


# ============================================================================
# Markdown converter
# ============================================================================


def _convert_markdown(source: Path, out_path: Path):
    """Convert Markdown to PDF via HTML intermediate."""
    text = source.read_text(encoding="utf-8", errors="replace")

    try:
        import markdown

        html = markdown.markdown(text, extensions=["tables", "fenced_code"])
    except ImportError:
        # Fallback: treat as plain text
        log.info("markdown library not installed, treating %s as plain text", source.name)
        return _convert_text(source, out_path)

    # Write temporary HTML, then convert via HTML converter
    import tempfile

    with tempfile.NamedTemporaryFile(
        mode="w",
        suffix=".html",
        delete=False,
        encoding="utf-8",
    ) as tmp:
        tmp.write(f"<html><body>{html}</body></html>")
        tmp_path = Path(tmp.name)

    try:
        _convert_html(tmp_path, out_path)
    finally:
        with contextlib.suppress(OSError):
            tmp_path.unlink()


# ============================================================================
# Plain text converter
# ============================================================================

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)(?:\s+#+)?\s*$")
_MD_BOLD_LINE_RE = re.compile(r"^\*\*(.+)\*\*$")
_HEADING_SIZES = {1: 22, 2: 18, 3: 15, 4: 13, 5: 12, 6: 11}


def _convert_text(source: Path, out_path: Path):
    """Convert plain text (with optional Markdown headings) to PDF.

    Recognises ``# heading`` syntax and ``**bold-only**`` lines,
    rendering them at larger font sizes so that downstream PDF
    parsers can detect heading structure via font-size distribution.
    """
    text = source.read_text(encoding="utf-8", errors="replace")

    pdf, font_name = _make_pdf()
    pdf.add_page()

    base_size = 10
    line_h = 5

    for line in text.split("\n"):
        stripped = line.strip()
        if not stripped:
            pdf.ln(line_h)
            pdf.set_x(pdf.l_margin)
            continue

        # --- Markdown heading: # ... ######
        hm = _MD_HEADING_RE.match(stripped)
        if hm:
            level = len(hm.group(1))
            title = hm.group(2).strip()
            # Strip inline bold markers from title text
            title = re.sub(r"\*\*(.+?)\*\*", r"\1", title)
            size = _HEADING_SIZES.get(level, base_size)
            h = max(size * 0.55, line_h)
            pdf.ln(3)
            _set_font(pdf, font_name, style="B", size=size)
            pdf.multi_cell(0, h, title)
            pdf.ln(2)
            _set_font(pdf, font_name, size=base_size)
            pdf.set_x(pdf.l_margin)
            continue

        # --- Bold-only line (potential sub-heading): **text**
        bm = _MD_BOLD_LINE_RE.match(stripped)
        if bm:
            title = bm.group(1).strip()
            _set_font(pdf, font_name, style="B", size=13)
            pdf.multi_cell(0, 7, title)
            _set_font(pdf, font_name, size=base_size)
            pdf.set_x(pdf.l_margin)
            continue

        # --- Regular line
        pdf.multi_cell(0, line_h, line)
        # fpdf2 bug: multi_cell(0, h, ...) can leave the cursor at the
        # right page edge when the last rendered fragment fills the line
        # exactly (trailing spaces, empty strings).  Reset X to the left
        # margin so the next call always has room to render.
        pdf.set_x(pdf.l_margin)

    pdf.output(str(out_path))
